import os
import pandas as pd
from pptx import Presentation
import win32com.client


def ask_if_missing(value, field_name, fmt_hint=None, default=None):
    """Ask user for input if value missing, otherwise return existing or default"""
    if pd.isnull(value) or str(value).strip() == "":
        if default is not None:
            return default
        if fmt_hint:
            user_input = input(f"Enter {field_name} ({fmt_hint}): ")
        else:
            user_input = input(f"Enter {field_name}: ")
        
        # If user enters just a space, treat as empty
        if user_input == " ":
            return ""
        return user_input.strip()
    return str(value).strip()


def pptx_to_pdf(pptx_path, pdf_path):
    """Convert PPTX to PDF using PowerPoint (Office 365 via COM)."""
    powerpoint = win32com.client.Dispatch("PowerPoint.Application")
    powerpoint.Visible = 1
    try:
        presentation = powerpoint.Presentations.Open(os.path.abspath(pptx_path))
        presentation.SaveAs(os.path.abspath(pdf_path), 32)  # 32 = PDF format
        presentation.Close()
    finally:
        powerpoint.Quit()


def replace_placeholders(slide, replacements):
    """Replace placeholders in a slide while preserving formatting."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    for key, value in replacements.items():
                        if key in run.text:
                            run.text = run.text.replace(key, value)


def normalize_validation(value):
    """Ensure validation ends with 'Years'"""
    if not value:
        return "3 Years"
    value = str(value).strip()
    if value.lower().endswith("years"):
        return value
    if value.isdigit():
        return f"{value} Years"
    return value


def clean_template_no(template_no, last_template):
    """Clean template number (Excel may return float like 1.0)."""
    if pd.isnull(template_no) or str(template_no).strip() == "":
        return last_template if last_template else 1
    try:
        return int(float(template_no))  # handles 1.0 → 1
    except ValueError:
        return last_template if last_template else 1


def get_template_file(template_no, course, last_template):
    """Decide which template file to use"""
    # Special case: LEVEL 1 HEALTH AND SAFETY (case-insensitive, flexible)
    if "level 1 health and safety" in course.lower():
        # If template number is provided, use it; otherwise default to template 2
        if pd.isnull(template_no) or str(template_no).strip() == "":
            template_no = 2
        else:
            template_no = clean_template_no(template_no, last_template)
    else:
        template_no = clean_template_no(template_no, last_template)

    template_file = f"templates/certificate_template_{template_no}.pptx"

    if not os.path.exists(template_file):
        raise FileNotFoundError(f"❌ Template file not found: {template_file}")

    return template_file, template_no


def sanitize_filename(text):
    """Remove forbidden characters from Windows filenames"""
    return "".join(c for c in str(text) if c not in r'\/:*?"<>|').strip()


def generate_certificates(excel_file):
    # Load Excel
    df = pd.read_excel(excel_file)

    if df.empty:
        print("❌ Excel file is empty.")
        return

    # Track last known values
    last_course = None
    last_date = None
    last_tutor = None
    last_validation = None
    last_template = None
    
    # Track if user has been asked for input (to avoid asking multiple times)
    user_asked_course = False
    user_asked_tutor = False
    user_asked_date = False

    certificates = []

    for _, row in df.iterrows():
        name = str(row.get("Name", "")).strip()
        if not name:
            continue

        # Course
        course = row.get("Course", "")
        if pd.notnull(course) and str(course).strip():
            course = str(course).strip()
            last_course = course
        elif last_course:
            course = last_course
        elif not user_asked_course:
            course = ask_if_missing("", "Course")
            if course:
                last_course = course
            user_asked_course = True
        else:
            course = last_course if last_course else ""

        # Tutor
        tutor = row.get("Tutor", "")
        if pd.notnull(tutor) and str(tutor).strip():
            tutor = str(tutor).strip()
            last_tutor = tutor
        elif last_tutor is not None:  # Use last_tutor if it exists (could be empty string)
            tutor = last_tutor
        elif not user_asked_tutor:
            # Ask for tutor input once, but allow it to be empty
            tutor_input = ask_if_missing("", "Tutor")
            if tutor_input:
                tutor = tutor_input
                last_tutor = tutor
            else:
                tutor = ""
                last_tutor = ""  # Remember that user chose empty
            user_asked_tutor = True
        else:
            tutor = last_tutor if last_tutor is not None else ""

        # Validation
        validation = row.get("Validation", "")
        validation = normalize_validation(
            validation if pd.notnull(validation) and str(validation).strip() else last_validation
        )
        last_validation = validation

        # Date (Excel dd-mm-yyyy parse, keep PDF long format)
        date_raw = row.get("Date", "")
        parsed_date = None

        if pd.notnull(date_raw) and str(date_raw).strip() != "":
            if isinstance(date_raw, (int, float)):  # Excel serial date
                try:
                    parsed_date = pd.to_datetime("1899-12-30") + pd.to_timedelta(int(date_raw), unit="D")
                except Exception:
                    parsed_date = pd.to_datetime(str(date_raw), dayfirst=True, errors="coerce")
            else:  # normal string or datetime
                parsed_date = pd.to_datetime(date_raw, dayfirst=True, errors="coerce")

        if (parsed_date is pd.NaT or parsed_date is None) and last_date:
            parsed_date = last_date

        if (parsed_date is pd.NaT or parsed_date is None) and not user_asked_date:
            date_input = ask_if_missing("", "Date", "dd-mm-yyyy")
            parsed_date = pd.to_datetime(date_input, dayfirst=True, errors="coerce")
            user_asked_date = True

        if parsed_date is pd.NaT or parsed_date is None:
            if last_date:
                parsed_date = last_date
            else:
                raise ValueError(f"❌ Could not parse date for row with Name='{name}'. Please ensure Date is in dd-mm-yyyy format.")

        last_date = parsed_date

        # ✅ Inside certificate/PDF → long format (25 September 2025)
        cert_date = parsed_date.strftime("%d %B %Y")
        # Folder name → ISO format (2025-09-25)
        folder_date = parsed_date.strftime("%Y-%m-%d")

        # Template
        template_no = row.get("Template", "")
        
        # Use last template if current row has no template specified
        if pd.isnull(template_no) or str(template_no).strip() == "":
            template_no = last_template
        
        template_file, current_template = get_template_file(template_no, course, last_template)
        last_template = current_template

        certificates.append({
            "name": name.upper(),
            "course": course.upper(),
            "tutor": tutor,
            "validation": validation,
            "cert_date": cert_date,
            "folder_date": folder_date,
            "template_file": template_file,
            "template_no": current_template
        })

    total = len(certificates)
    if total == 0:
        print("❌ No valid rows found in Excel.")
        return

    for idx, cert in enumerate(certificates, start=1):
        name = sanitize_filename(cert["name"])
        course = sanitize_filename(cert["course"])
        tutor = cert["tutor"]
        validation = cert["validation"]
        cert_date = cert["cert_date"]
        folder_date = cert["folder_date"]
        template_file = cert["template_file"]

        # Folder: {course}/{yyyy-mm-dd}/
        course_folder = os.path.join(course, folder_date)
        pptx_folder = os.path.join(course_folder, "PPTX")
        os.makedirs(pptx_folder, exist_ok=True)

        pptx_file = os.path.join(pptx_folder, f"Certificate on {course} - {name}.pptx")
        pdf_file = os.path.join(course_folder, f"Certificate on {course} - {name}.pdf")

        print(f"\nCreating certificate for {name} ({idx}/{total}) | Course: {course} | Date: {cert_date} | Template: {template_file}")

        # Load PPTX template
        prs = Presentation(template_file)
        slide = prs.slides[0]

        # Replace placeholders
        replacements = {
            "{{NAME}}": name,
            "{{COURSE}}": course,
            "{{DATE}}": cert_date,
            "{{TUTOR}}": tutor,
            "{{VALIDATION}}": validation
        }
        replace_placeholders(slide, replacements)

        # Save PPTX
        prs.save(pptx_file)

        # Convert to PDF
        pptx_to_pdf(pptx_file, pdf_file)

        print(f"✅ Created: {pdf_file}")

    print(f"\n🎉 All {total} certificates generated successfully.")


if __name__ == "__main__":
    generate_certificates("students.xlsx")
